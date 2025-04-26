import re
import json
import traceback
from typing import List, Dict
from pprint import pprint, pformat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from infrastructure.ppt_reader import PPTReader
from infrastructure.open_microsoft_document import IOpenAndUpdateDocument
from infrastructure.generic_logger import GenericLogger
from domain.llm_utils import LLMUtils
from domain.queue import MetadataPpt
from domain.worker_class import Worker

class OpenPPTDocument(IOpenAndUpdateDocument): 
    def __init__(self, document_path: str, 
                 worker: Worker, 
                 paragraph_start_min_word_numbers: int,
                 paragraph_start_min_word_length: int,
                 slides_to_skip: List, slides_to_keep: List,\
                 logger: GenericLogger, llm_utils: LLMUtils):
        self.logger = logger
        self.document =  Presentation(document_path)
        self.llm_utils = llm_utils
        self.document_path = document_path
        self.slides_to_skip = slides_to_skip
        self.slides_to_keep = slides_to_keep
        self.paragraph_start_min_word_numbers = paragraph_start_min_word_numbers
        self.paragraph_start_min_word_length = paragraph_start_min_word_length

        super().__init__(document_path, worker, 
                         paragraph_start_min_word_numbers, paragraph_start_min_word_length, 
                         logger)
    
    def __print_slide_keep_skip_info(self, keep_skip_info: str) -> None:
        self.logger.log_info(keep_skip_info)

    def __append(self, shape_descriptions: List, shape: Dict) -> None:
        if shape is not None:
            shape_descriptions.append(shape)
        else:
            self.logger.log_debug("None shape was not added")

    def __ppt_to_json(self):

        # This method returns each slide in the following data structure
        # shape_descriptions = [
        #  {
        #    'json': {
        #            "text": [shape.text_frame.text], // We need array for groups
        #            "pointers": [shape.text_frame],
        #            "is_title": True / False
        #     },
        #    'raw_text': str(new_json["shape"]["text"]) if text is None else text
        #    "slide_number": slide_number,
        #  }
        #]

        for slide_idx, slide in enumerate(self.document.slides):

            slide_number: int = slide_idx + 1
            if slide_number in self.slides_to_skip:
                self.__print_slide_keep_skip_info(f"Skipped slide number {slide_number} as per request.")
                continue

            if self.slides_to_keep is not None and len(self.slides_to_keep) > 0:
                if slide_number not in self.slides_to_keep:
                    continue
                else:
                    self.__print_slide_keep_skip_info(f"Keep slide number {slide_number} as per request.")               

            if slide.element.get('show', '1' == '0'):
                self.__print_slide_keep_skip_info(f"Skipped hidden slide number {slide_number}.")
                continue
 
            self.logger.log_info(f"Analyzing slide number {slide_number}")
 
            shape_descriptions: list = [] 
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self.__append(shape_descriptions, PPTReader.get_text_box_info(slide_number, shape, self))
 
                elif shape.has_table: 
                    table = shape.table
                    pointer_list: List = [table]
                    table_md_str: str = ""
                    first_row: bool = True
                    for row in table.rows:
                        table_md_str += "\n|"
                        for cell in row.cells:
                            text: str = ""
                            if cell.text_frame is not None:
                              text = cell.text_frame.text.replace("\n", " ")
                            table_md_str += text + "|"
                        if first_row: table_md_str += "\n" + ("-" * max(3, len(table_md_str)))
                        first_row = False
                        
                    self.__append(shape_descriptions, PPTReader.get_table_info(slide_number, shape, table_md_str, pointer_list, self))

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: 
                    self.__append(shape_descriptions, PPTReader.get_group_info(slide_number, shape, self))
               
                else: 
                    self.__append(shape_descriptions, PPTReader.get_shape_type_info(slide_number, shape, self))
 
            title_value: str = None
            shape_title: Dict = None
            title_found: bool = False
            if hasattr(slide.shapes, "title") and hasattr(slide.shapes.title, "text") and \
              slide.shapes.title.text is not None and self.is_paragraph(slide.shapes.title.text) > 0:
                title_value = slide.shapes.title.text
                for shape_description in shape_descriptions:
                    if shape_description["raw_text"] == title_value:
                        shape_title = shape_description
                        shape_title["json"]["is_title"] = True
                        title_found = True

                        self.logger.log_trace(f"Recognized title from shape_description: {pformat(shape_title)}, title_value = {title_value}")
                if (not title_found):
                    shape_description: Dict = PPTReader.create_title(slide_number, slide.shapes.title, self)
                    shape_title = shape_description
                    shape_descriptions.append(shape_description)
                    title_found = True
                    self.logger.log_trace(f"Creating default title: {pformat(shape_description)}, slide.shapes.title = {slide.shapes.title}")

            sorted_shapes: List = PPTReader.get_sorted_shapes_by_pos_y(shape_descriptions)

            context: str = ""
            for shape_description in sorted_shapes:
                text = shape_description['raw_text']
                if self.is_paragraph(text):
                    if shape_description['json']['is_title'] == True:
                        context += "# "
                        context += text + "\n"

            for shape_description in sorted_shapes:
                if isinstance(shape_description['raw_text'], str):
                    text: str = shape_description['raw_text']
                    text = re.sub(r"^\s*$", "", text)
                    if len(text) > 0:
                        request_type: str = LLMUtils.DEFAULT_REQUEST
                        if shape_description['json']['type'] == str(MSO_SHAPE_TYPE.TABLE):
                            request_type = LLMUtils.TABLE_REQUEST
                        if shape_description['json']['is_title'] == True:
                            self.logger.log_trace(f"Adding heading {text}")
                            request_type = LLMUtils.HEADING_REQUEST

                        self.logger.log_trace(f"Populating requests: shape_description = {pformat(shape_description)}")
                        self.worker.add_work_element(MetadataPpt(shape_description["json"]["pointers"], \
                                                                context, text, \
                                                                request_type, self.logger))   
            
    def process(self):
        self.__ppt_to_json()
        self.worker.process_all()

        