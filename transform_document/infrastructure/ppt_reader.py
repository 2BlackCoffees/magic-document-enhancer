from typing import Dict, List
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint, pformat
from pathlib import Path
from infrastructure.open_microsoft_document import IOpenAndUpdateDocument
class PPTReader:

    @staticmethod  
    def __check_recursively_for_text(shape, json_shape, open_and_update_document: IOpenAndUpdateDocument):
        for cur_shape in shape.shapes:
            if cur_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                PPTReader.__check_recursively_for_text(cur_shape, json_shape, open_and_update_document)
            else:
                if hasattr(cur_shape, "text") and open_and_update_document.is_paragraph(cur_shape.text):
                    json_shape["text"].append(cur_shape.text)
                    json_shape["pointers"].append(cur_shape)
        
        return  json_shape
    
    @staticmethod  
    def _get_shape_infos(shape: Dict, type: MSO_SHAPE_TYPE, open_and_update_document: IOpenAndUpdateDocument) -> Dict:

        if hasattr(shape, 'text_frame'):
            text_frame: any = shape.text_frame
            if open_and_update_document.is_paragraph(text_frame.text):
                return  {
                        'y': shape.top, 'x': shape.left, 
                        "text": [str(text_frame.text)],
                        "pointers": [text_frame],
                        "is_title": "False",
                        "type": str(type)
                }

        return None
    
    @staticmethod  
    def _get_shape_group_infos(shape: Dict, type: MSO_SHAPE_TYPE, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        text: List = []
        pointers: List = []
        if hasattr(shape, 'text_frame') and open_and_update_document.is_paragraph(shape.text_frame.text):
            text = [str(shape.text_frame.text)]
            pointers = [shape.text_frame]

            return  {
                    'y': shape.top, 'x': shape.left, 
                    "text": text,
                    "pointers": pointers,
                    "is_title": "False",
                    "type": str(type)
            }

        return None
    
    @staticmethod  
    def _get_title_infos(shape: Dict, type: MSO_SHAPE_TYPE, open_and_update_document: IOpenAndUpdateDocument) -> Dict:

        if hasattr(shape, 'text') and open_and_update_document.is_paragraph(str(shape.text)):
            return  {
                    'y': shape.top, 'x': shape.left, 
                    "text": [str(shape.text)],
                    "pointers": [shape],
                    "is_title": "True",
                    "type": str(type)
            }

        return None

    @staticmethod  
    def _get_shape_table_infos(shape: Dict, pointer_list: List, table_str: str, type: MSO_SHAPE_TYPE) -> Dict:

        return  {
                'y': shape.top, 'x': shape.left, 
                "text": table_str,
                "pointers": pointer_list,
                "is_title": "False",
                "type": str(type)
        }

    @staticmethod  
    def __encapsulate_shape(slide_number: int, json_shape: Dict, open_and_update_document: IOpenAndUpdateDocument, raw_text: str = None) -> Dict:
        if json_shape is not None:
            new_json = json_shape.copy()
            text = " ".join(new_json["text"]) if raw_text is None else raw_text
            if open_and_update_document.is_paragraph(text):
                return {
                    'y': json_shape['y'], 
                    'json': new_json,
                    'raw_text': text,
                    'slide_number': slide_number
                }
        
        return None

    @staticmethod  
    def get_text_box_info(slide_number: int, shape: Dict, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        json_shape: Dict = PPTReader._get_shape_infos(shape, MSO_SHAPE_TYPE.TEXT_BOX, open_and_update_document)
        return PPTReader.__encapsulate_shape(slide_number, json_shape, open_and_update_document)
    
    @staticmethod  
    def get_group_info(slide_number: int, shape: Dict, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        json_shape: Dict = PPTReader._get_shape_group_infos(shape, MSO_SHAPE_TYPE.GROUP, open_and_update_document)
        if json_shape is not None:
            json_shape = PPTReader.__check_recursively_for_text(shape, json_shape, open_and_update_document)
        return PPTReader.__encapsulate_shape(slide_number, json_shape, open_and_update_document)
    
    @staticmethod  
    def get_table_info(slide_number: int, shape: Dict, table_str: str, pointer_list: List, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        json_shape: Dict = PPTReader._get_shape_table_infos(shape, pointer_list, table_str, MSO_SHAPE_TYPE.TABLE)
        return PPTReader.__encapsulate_shape(slide_number, json_shape, open_and_update_document, table_str)

    @staticmethod  
    def get_shape_type_info(slide_number: int, shape: Dict, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        json_shape: Dict = PPTReader._get_shape_infos(shape, MSO_SHAPE_TYPE.MIXED, open_and_update_document)
        return PPTReader.__encapsulate_shape(slide_number, json_shape, open_and_update_document)

    @staticmethod  
    def create_title(slide_number: int, shape: Dict, open_and_update_document: IOpenAndUpdateDocument) -> Dict:
        json_shape: Dict = PPTReader._get_title_infos(shape, MSO_SHAPE_TYPE.TEXT_BOX, open_and_update_document)
        return PPTReader.__encapsulate_shape(slide_number, json_shape, open_and_update_document)

    @staticmethod  
    def get_sorted_shapes_by_pos_y(shapes: List) -> List:
        return sorted(shapes, key = lambda shape_dict: shape_dict['y'])

